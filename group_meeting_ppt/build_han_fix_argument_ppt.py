from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/han/Ai_ws/Study/vrx_ws")
OUT = ROOT / "group_meeting_ppt"
ASSETS = OUT / "assets"
PPTX_OUT = OUT / "han_usv_controller_本条对话修改论述流程汇报.pptx"
SCRIPT_OUT = OUT / "han_usv_controller_本条对话修改论述流程汇报讲稿.md"
FLOW_OUT = OUT / "han_usv_controller_本条对话修改流程图.mmd"

FONT = "Noto Sans CJK SC"
C = {
    "bg": RGBColor(244, 247, 249),
    "white": RGBColor(255, 255, 255),
    "navy": RGBColor(23, 50, 77),
    "dark": RGBColor(16, 37, 52),
    "text": RGBColor(37, 49, 59),
    "muted": RGBColor(96, 112, 125),
    "line": RGBColor(213, 222, 229),
    "pale": RGBColor(232, 238, 243),
    "teal": RGBColor(36, 127, 120),
    "blue": RGBColor(44, 120, 165),
    "green": RGBColor(46, 139, 98),
    "orange": RGBColor(214, 122, 45),
    "red": RGBColor(183, 76, 70),
    "cyan": RGBColor(26, 175, 193),
}


def set_text(tf, text, size=14, color="text", bold=False, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.color.rgb = C[color] if isinstance(color, str) else color
    p.font.bold = bold
    if align:
        p.alignment = align
    return p


def add_text(slide, text, x, y, w, h, size=14, color="text", bold=False, align=None, valign=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    if valign:
        box.text_frame.vertical_anchor = valign
    set_text(box.text_frame, text, size=size, color=color, bold=bold, align=align)
    return box


def add_rect(slide, x, y, w, h, fill="white", line="line", radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C[fill] if isinstance(fill, str) else fill
    shp.line.color.rgb = C[line] if isinstance(line, str) else line
    shp.line.width = Pt(1)
    return shp


def add_card(slide, x, y, w, h, title, body, accent="teal", idx=None):
    add_rect(slide, x, y, w, h)
    add_rect(slide, x, y, 0.08, h, fill=accent, line=accent)
    if idx:
        add_text(slide, idx, x + 0.22, y + 0.18, 0.62, 0.35, size=18, color=accent, bold=True)
        title_x = x + 0.86
        title_w = w - 1.06
    else:
        title_x = x + 0.24
        title_w = w - 0.48
    add_text(slide, title, title_x, y + 0.2, title_w, 0.34, size=14.5, color="navy", bold=True)
    add_text(slide, body, x + 0.24, y + 0.7, w - 0.48, h - 0.86, size=10.9, color="text")


def add_base(prs, title, section, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C["bg"]
    add_rect(slide, 0, 0, 13.333, 0.12, fill="teal", line="teal")
    add_text(slide, title, 0.62, 0.34, 10.6, 0.48, size=23.5, color="navy", bold=True)
    add_text(slide, section, 10.65, 0.42, 2.0, 0.24, size=9.5, color="muted", align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.62), Inches(0.94), Inches(12.70), Inches(0.94))
    line.line.color.rgb = C["line"]
    line.line.width = Pt(1)
    add_text(slide, "VRX WAM-V 自主航行系统 | han_usv_controller 本条对话修改", 0.62, 7.06, 6.3, 0.24, size=8.5, color="muted")
    add_text(slide, str(page), 12.35, 7.06, 0.35, 0.24, size=9, color="muted", align=PP_ALIGN.RIGHT)
    return slide


def add_arrow(slide, x1, y1, x2, y2, color="line", width=1.4):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = C[color]
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def add_flow_box(slide, x, y, w, h, title, subtitle, color="teal"):
    add_rect(slide, x, y, w, h)
    add_rect(slide, x, y, w, 0.11, fill=color, line=color)
    add_text(slide, title, x + 0.1, y + 0.24, w - 0.2, 0.28, size=12.2, color="navy", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, x + 0.12, y + 0.62, w - 0.24, h - 0.72, size=8.8, color="muted", align=PP_ALIGN.CENTER)


def add_metric(slide, x, y, value, label, color="teal"):
    add_rect(slide, x, y, 2.35, 1.0)
    add_text(slide, value, x + 0.12, y + 0.18, 2.11, 0.34, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.16, y + 0.62, 2.03, 0.24, size=9.4, color="muted", align=PP_ALIGN.CENTER)


def image_or_placeholder(slide, path, x, y, w, h, label):
    if path.exists():
        picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y))
        target_ratio = w / h
        image_ratio = picture.width / picture.height
        if image_ratio > target_ratio:
            crop = (1 - target_ratio / image_ratio) / 2
            picture.crop_left = crop
            picture.crop_right = crop
        else:
            crop = (1 - image_ratio / target_ratio) / 2
            picture.crop_top = crop
            picture.crop_bottom = crop
        picture.left = Inches(x)
        picture.top = Inches(y)
        picture.width = Inches(w)
        picture.height = Inches(h)
        return picture
    else:
        add_rect(slide, x, y, w, h, fill="pale")
        add_text(slide, label, x + 0.2, y + h / 2 - 0.15, w - 0.4, 0.3, size=12, color="muted", align=PP_ALIGN.CENTER)


def add_image_caption(slide, path, x, y, w, h, title, accent="teal"):
    image_or_placeholder(slide, path, x, y, w, h, title)
    band = add_rect(slide, x, y + h - 0.48, w, 0.48, fill="dark", line="dark")
    band.fill.transparency = 12
    add_rect(slide, x, y + h - 0.48, 0.08, 0.48, fill=accent, line=accent)
    add_text(slide, title, x + 0.2, y + h - 0.35, w - 0.35, 0.2, size=10.2, color="white", bold=True)


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.author = "HanW"
    prs.core_properties.title = "han_usv_controller 本条对话修改论述流程汇报"
    prs.core_properties.subject = "VRX WAM-V han_usv_controller 改造流程与验收"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_or_placeholder(slide, ASSETS / "sim_title.png", 0, 0, 13.333, 7.5, "VRX 仿真标题图")
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = C["dark"]
    overlay.fill.transparency = 30
    overlay.line.fill.background()
    add_rect(slide, 0.7, 1.12, 0.08, 3.65, fill="orange", line="orange")
    add_text(slide, "han_usv_controller 本条对话修改\n论述流程汇报", 1.02, 1.26, 9.8, 1.45, size=30, color="white", bold=True)
    add_text(slide, "从“能跑”走向“可观测、可复现、可验收”的 VRX WAM-V 自主航行闭环", 1.04, 3.03, 9.9, 0.42, size=16.5, color=RGBColor(232, 241, 245))
    tags = ["EKF 定位", "滚动地图", "State Lattice", "COLREGs", "自动评测"]
    tx = 1.04
    for tag in tags:
        add_rect(slide, tx, 3.77, 1.22 if len(tag) < 6 else 1.42, 0.38, fill="dark", line="white")
        add_text(slide, tag, tx, 3.84, 1.22 if len(tag) < 6 else 1.42, 0.18, size=9.6, color="white", align=PP_ALIGN.CENTER)
        tx += (1.22 if len(tag) < 6 else 1.42) + 0.16
    add_text(slide, "汇报主线：问题牵引 -> 系统改造 -> 启动流程 -> 控制流程 -> 验收证据 -> 后续边界", 1.04, 4.52, 9.6, 0.28, size=11.4, color=RGBColor(232, 241, 245))
    add_text(slide, "汇报人：HanW | 2026.07.20", 1.04, 6.72, 5.0, 0.28, size=11, color="white")
    add_text(slide, "论述 PPT", 10.22, 6.66, 2.1, 0.32, size=11.5, color="white", bold=True, align=PP_ALIGN.RIGHT)

    slide = add_base(prs, "一、为什么要改：把 VRX 练习变成工程闭环", "问题与论点", 2)
    add_image_caption(slide, ASSETS / "sim_card.png", 0.75, 1.22, 12.0, 2.3, "真实 VRX Wayfinding 场景：到点只是结果，过程仍需被解释", "orange")
    add_card(slide, 0.75, 3.78, 3.72, 2.05, "原始痛点", "定位源、障碍记忆和动态船预测不透明，成功或失败很难复盘。", "red", "01")
    add_card(slide, 4.82, 3.78, 3.72, 2.05, "本轮论点", "每个决策有可视化、每次运行有指标、每类风险有独立验收。", "teal", "02")
    add_card(slide, 8.89, 3.78, 3.72, 2.05, "工程标准", "同一推荐入口整合算法、场景、RViz 与评测，形成一致实验链。", "blue", "03")
    add_text(slide, "结论：han_usv_controller 把“能跑”提升为可学习、可诊断、可回归。", 0.82, 6.12, 11.8, 0.42, size=15.2, color="navy", bold=True, align=PP_ALIGN.CENTER)

    slide = add_base(prs, "二、从这条对话开始的修改范围", "改造总览", 3)
    add_image_caption(slide, ASSETS / "sim_card.png", 4.52, 1.5, 4.3, 3.2, "场景、算法与验收围绕同一艘 WAM-V 闭环", "teal")
    add_card(slide, 0.72, 1.28, 3.42, 2.25, "感知与定位", "EKF / 健康回退\n点云滤水 / 浮标聚类\n动态目标跟踪", "teal")
    add_card(slide, 0.72, 3.82, 3.42, 2.25, "地图与规划", "100 m 滚动栅格\nDubins State Lattice A*\n在线重规划 / fallback 监控", "blue")
    add_card(slide, 9.18, 1.28, 3.42, 2.25, "控制与规则", "ILOS + PID / 曲率限速\n制动与倒车恢复\nCPA/TCPA / COLREGs", "green")
    add_card(slide, 9.18, 3.82, 3.42, 2.25, "场景与验收", "12 枚浮标 / 2 艘动态船\n中文 RViz / 仿真独占锁\n自动评测硬门槛", "orange")
    add_text(slide, "定位→建图→规划→控制→规则→可视化→验收", 4.38, 5.05, 4.58, 0.54, size=12.6, color="navy", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "范围边界：未改旧版 autonomous_controller.py / launch_wayfinding.sh；推荐演示入口集中到 han_usv_controller。", 0.82, 6.44, 11.8, 0.28, size=11.6, color="muted", align=PP_ALIGN.CENTER)

    slide = add_base(prs, "三、演示启动流程：一条命令串起所有层", "运行流程", 4)
    steps = [
        ("清理环境", "停止旧 Gazebo\n检查残留进程", "red"),
        ("首次构建", "colcon build\nsource install", "orange"),
        ("主场景启动", "colregs_learning\n.launch.py\n推荐入口", "teal"),
        ("基础仿真", "simulation\n.launch.py\nWayfinding + WAM-V", "blue"),
        ("算法节点", "EKF + Map + Lattice\nILOS/PID + COLREGs", "green"),
        ("可视化/验收", "RViz 观察\nstatus + 自动评测", "cyan"),
    ]
    x = 0.66
    for i, (title, subtitle, color) in enumerate(steps):
        add_flow_box(slide, x, 1.42, 1.75, 1.15, title, subtitle, color)
        if i < len(steps) - 1:
            add_arrow(slide, x + 1.77, 1.98, x + 2.15, 1.98, color="muted")
        x += 2.05
    add_rect(slide, 0.78, 3.18, 7.25, 2.55)
    add_text(slide, "启动链核心", 1.02, 3.43, 2.0, 0.32, size=15, color="navy", bold=True)
    chain = (
        "colregs_learning.launch.py  ->  simulation.launch.py  ->  官方 Wayfinding / WAM-V\n"
        "  ->  robot_localization ekf_node  ->  autonomous_usv\n"
        "  ->  滚动占据栅格 / State Lattice / ILOS + PID / COLREGs Marker\n"
        "  ->  12 枚红绿浮标 + 2 艘动态目标船 + pointcloud.rviz"
    )
    add_text(slide, chain, 1.02, 3.94, 6.72, 1.18, size=10.8, color="text")
    add_image_caption(slide, ASSETS / "sim_card.png", 8.35, 3.18, 4.25, 2.55, "启动后：Gazebo 场景、WAM-V 与动态目标同步加载", "orange")
    add_text(slide, "关键约束：同一时刻只运行一套 VRX 仿真；重复启动会被独占锁拒绝。", 1.02, 5.92, 11.2, 0.3, size=11, color="muted")

    slide = add_base(prs, "四、控制流程：静态避障与动态会遇分层接入", "算法流程", 5)
    left_steps = [
        ("传感器输入", "GNSS / IMU / Scan / PointCloud", "teal"),
        ("状态估计", "robot_localization EKF\nPlanarEKF 兜底", "blue"),
        ("滚动地图", "100 m 栅格\n占用/自由/衰减", "green"),
        ("路径规划", "Dubins State Lattice A*\n40 m 子目标", "orange"),
        ("路径跟踪", "ILOS + 曲率限速\n艏摇前馈", "cyan"),
        ("执行输出", "安全过滤\n左右推进器", "red"),
    ]
    x = 0.62
    for i, (title, subtitle, color) in enumerate(left_steps):
        add_flow_box(slide, x, 1.35, 1.75, 1.0, title, subtitle, color)
        if i < len(left_steps) - 1:
            add_arrow(slide, x + 1.76, 1.85, x + 2.08, 1.85, color="muted")
        x += 2.04
    add_image_caption(slide, ASSETS / "lidar_card.png", 0.92, 3.15, 3.15, 2.35, "静态层：点云 → 栅格 → State Lattice", "blue")
    add_rect(slide, 4.3, 3.15, 4.72, 2.35)
    add_text(slide, "动态目标监督分支", 4.58, 3.4, 3.1, 0.3, size=14.5, color="navy", bold=True)
    add_text(slide, "2 艘目标船 → alpha-beta 跟踪 → CPA/TCPA → 会遇分类 → 右转偏置 + 降速", 4.58, 3.88, 4.16, 0.75, size=11.2, color="text")
    add_text(slide, "动态目标不写成静态墙；原始激光仍进入近障碍安全层。", 4.58, 4.77, 4.12, 0.38, size=10.3, color="muted")
    add_image_caption(slide, ASSETS / "path_card.png", 9.25, 3.15, 3.15, 2.35, "执行层：ILOS 贴线与实际轨迹", "green")
    add_text(slide, "这保证动态船不会变成地图里的“虚假墙”，也不会绕开近距离安全检查。", 0.82, 6.13, 11.7, 0.34, size=12.2, color="navy", bold=True, align=PP_ALIGN.CENTER)

    slide = add_base(prs, "五、可视化结果：让决策过程能被看见", "场景与 RViz", 6)
    image_or_placeholder(slide, ASSETS / "path_card.png", 0.75, 1.24, 5.55, 3.7, "路径跟踪截图")
    image_or_placeholder(slide, ASSETS / "lidar_card.png", 7.03, 1.24, 5.55, 3.7, "雷达与岸线截图")
    add_card(slide, 0.75, 5.28, 3.72, 1.18, "路径/轨迹/XTE", "青色规划路径、黄色实际轨迹、品红横向误差，把“船为什么这样走”变成可检查证据。", "blue")
    add_card(slide, 4.82, 5.28, 3.72, 1.18, "浮标候选", "黄色圆环来自真实三维激光聚类；定点删除 Marker，消除周期性闪烁。", "orange")
    add_card(slide, 8.89, 5.28, 3.72, 1.18, "动态船/COLREGs", "橙色方框显示动态船，风险目标变红，并显示 TCPA、DCPA、会遇类型。", "teal")

    slide = add_base(prs, "六、自动验收：用指标证明不是只会演示", "评测证据", 7)
    add_metric(slide, 0.78, 1.28, "100%", "官方任务完成率", "green")
    add_metric(slide, 3.08, 1.28, "0", "碰撞次数", "green")
    add_metric(slide, 5.38, 1.28, "1.076 m", "官方平均误差", "teal")
    add_metric(slide, 7.68, 1.28, "5.952 m", "最小净空", "blue")
    add_metric(slide, 9.98, 1.28, "0", "State Lattice fallback", "green")
    add_metric(slide, 0.78, 2.58, "2", "最大动态跟踪数", "teal")
    add_metric(slide, 3.08, 2.58, "189", "COLREGs 激活样本", "orange")
    add_metric(slide, 5.38, 2.58, "11.244 deg/s", "最大艏摇角速度", "blue")
    add_metric(slide, 7.68, 2.58, "6.151 ms", "最大规划耗时", "teal")
    add_metric(slide, 9.98, 2.58, "172 passed", "源码/配置/模型测试", "green")
    add_rect(slide, 0.9, 4.34, 7.3, 1.55)
    add_text(slide, "验收命令", 1.16, 4.58, 1.6, 0.28, size=14.5, color="navy", bold=True)
    add_text(slide, "bash han_usv_controller/scripts/run_colregs_learning.sh", 2.7, 4.58, 5.1, 0.28, size=12.2, color="teal", bold=True)
    add_text(slide, "硬门槛：三航点、0 碰撞、EKF 健康、动态目标=2、COLREGs>=20、fallback=0。", 1.16, 5.08, 6.7, 0.46, size=10.6, color="text")
    add_image_caption(slide, ASSETS / "path_card.png", 8.48, 4.34, 3.94, 1.55, "验收实景：规划路径与实际轨迹同时留证", "green")
    add_text(slide, "证据：han_usv_controller/evaluation/20260717T002427Z/aggregate.json", 1.16, 6.08, 10.9, 0.28, size=10.6, color="muted")

    slide = add_base(prs, "七、论述结论：这次改造解决了什么", "贡献与边界", 8)
    add_image_caption(slide, ASSETS / "sim_card.png", 0.75, 1.25, 3.72, 2.25, "体系化：场景与推荐入口统一", "teal")
    add_image_caption(slide, ASSETS / "lidar_card.png", 4.82, 1.25, 3.72, 2.25, "可解释：感知和避障过程可见", "blue")
    add_image_caption(slide, ASSETS / "path_card.png", 8.89, 1.25, 3.72, 2.25, "可回归：路径与轨迹可量化", "green")
    add_card(slide, 0.75, 3.72, 3.72, 2.12, "贡献 1：体系化", "launch、配置、模型、算法、脚本和测试各有边界，推荐入口统一。", "teal", "A")
    add_card(slide, 4.82, 3.72, 3.72, 2.12, "贡献 2：可解释", "路径、误差、目标、浮标与控制状态进入 RViz/status，便于复盘。", "blue", "B")
    add_card(slide, 8.89, 3.72, 3.72, 2.12, "贡献 3：可回归", "完成率、碰撞、净空、耗时与 COLREGs 激活都成为自动门槛。", "green", "C")
    add_text(slide, "当前边界：动态船来自仿真专用检测 Topic；激光不能识别红/绿语义；COLREGs 是教学监督层，不是认证避碰系统。", 0.82, 6.24, 11.75, 0.36, size=11.5, color="muted", align=PP_ALIGN.CENTER)

    slide = add_base(prs, "八、汇报时的推荐流程", "讲述路径", 9)
    agenda = [
        ("1 分钟", "先抛问题", "旧系统能跑，但难以解释、难以复现、难以验收。"),
        ("2 分钟", "讲改造结构", "定位、地图、规划、控制、规则、可视化、评测七层闭环。"),
        ("2 分钟", "展示启动链", "强调 colregs_learning.launch.py 是推荐入口，避免重复仿真。"),
        ("3 分钟", "讲算法流程", "静态障碍走地图和 State Lattice，动态目标走 CPA/TCPA/COLREGs 分支。"),
        ("2 分钟", "展示指标", "用 aggregate.json 的 PASS 指标收束论证。"),
    ]
    y = 1.28
    for time, title, body in agenda:
        add_rect(slide, 0.82, y, 1.2, 0.72, fill="teal", line="teal")
        add_text(slide, time, 0.87, y + 0.24, 1.1, 0.2, size=10.2, color="white", bold=True, align=PP_ALIGN.CENTER)
        add_card(slide, 2.15, y, 10.05, 0.72, title, body, "blue")
        y += 0.93
    add_text(slide, "收束句：这不是单个避障算法的展示，而是一套面向 VRX 学习和后续研究的可验证实验平台。", 0.95, 6.36, 11.3, 0.34, size=14.0, color="navy", bold=True, align=PP_ALIGN.CENTER)

    prs.save(PPTX_OUT)


def write_notes():
    SCRIPT_OUT.write_text(
        """# han_usv_controller 本条对话修改论述流程汇报讲稿

## 1. 标题页
这次汇报的核心不是单独展示一个控制器，而是说明 `han_usv_controller` 如何从一条能跑的 VRX 任务，变成一套可观测、可复现、可验收的自主航行实验平台。

## 2. 为什么要改
原始 Wayfinding 能证明船可以到点，但很难解释成功或失败的原因。本轮改造围绕工程闭环展开：定位、地图、规划、控制、动态目标、可视化和自动评测必须互相咬合。

## 3. 修改范围
本条对话的修改可以分成四类：感知定位、地图规划、控制规则、场景验收。注意说明旧版 `autonomous_controller.py` 和 `launch_wayfinding.sh` 不是本次推荐入口。

## 4. 启动流程
推荐启动命令固定为 `ros2 launch han_usv_controller colregs_learning.launch.py`。它内部包含官方 Wayfinding、WAM-V、EKF、控制器、RViz、12 枚浮标和 2 艘动态船。重复仿真会被独占锁拦住。

## 5. 控制流程
静态障碍走滚动占据栅格和 Dubins State Lattice，动态目标走 CPA/TCPA 和 COLREGs 监督。COLREGs 只给航向偏置和速度缩放，最终仍由 ILOS/PID 和安全层输出左右推进器。

## 6. 可视化结果
RViz 不只是好看，它把路径、轨迹、横向误差、动态目标、浮标候选、占据栅格都变成可检查证据。黄色浮标候选来自真实点云聚类，不再使用重复的红色通用障碍 Marker。

## 7. 自动验收
用 `run_colregs_learning.sh` 做闭环验收。重点讲 `100%` 完成、`0` 碰撞、`1.076 m` 官方平均误差、`2` 个动态目标、`189` 个 COLREGs 激活样本、`0` fallback。

## 8. 贡献与边界
贡献是体系化、可解释、可回归。边界也要说清楚：动态船检测是仿真专用 Topic，激光不能识别红绿语义，COLREGs 是教学监督层，不是认证系统。

## 9. 推荐讲述顺序
先讲问题，再讲结构，再讲启动流程和控制流程，最后用指标收束。结尾强调：这是一套面向 VRX 学习和后续研究的可验证实验平台。
""",
        encoding="utf-8",
    )

    FLOW_OUT.write_text(
        """flowchart TD
    A["本条对话修改需求"] --> B["功能拆分：定位 / 地图 / 规划 / 控制 / 规则 / 可视化 / 评测"]
    B --> C["实现 han_usv_controller ROS2 包"]
    C --> D["Launch 整合：colregs_learning.launch.py"]
    D --> E["仿真场景：Wayfinding + 12 枚浮标 + 2 艘动态船"]
    E --> F["控制链：EKF -> 滚动地图 -> State Lattice -> ILOS/PID"]
    E --> G["监督链：动态目标 -> CPA/TCPA -> COLREGs 偏置/降速"]
    F --> H["RViz：路径 / 轨迹 / XTE / 栅格 / 浮标候选"]
    G --> H
    H --> I["自动验收：run_colregs_learning.sh"]
    I --> J["证据：100% 完成、0 碰撞、189 个 COLREGs 激活样本、0 fallback"]
""",
        encoding="utf-8",
    )


if __name__ == "__main__":
    OUT.mkdir(parents=True, exist_ok=True)
    build_ppt()
    write_notes()
    print(PPTX_OUT)
    print(SCRIPT_OUT)
    print(FLOW_OUT)
