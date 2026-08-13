from pathlib import Path
from textwrap import dedent

from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/han/Ai_ws/Study/vrx_ws")
OUT = ROOT / "group_meeting_ppt"
ASSETS = OUT / "assets"
RENDERED = OUT / "rendered"

PPTX_OUT = OUT / "VRX_WAMV_自主航行项目方案扩展版.pptx"
NOTES_OUT = OUT / "VRX_WAMV_自主航行项目方案扩展版讲稿.md"
BG_OUT = ASSETS / "bit_green_watermark_bg.png"

SOURCE_SLIDES = [
    Path("/tmp/han-clipboard-9ae44990-2361-403d-ac03-100f631a86d9.png"),
    Path("/tmp/han-clipboard-9126e553-c0e7-45d4-b0f2-d60dfcb9f0e8.png"),
    Path("/tmp/han-clipboard-12384b0f-eab1-41fb-a48d-982fb4b6547b.png"),
    Path("/tmp/han-clipboard-c52b0693-eaa0-4681-b0b6-e2376453b2db.png"),
    Path("/tmp/han-clipboard-b826884d-6bcf-446f-963a-1bced1cf8a6b.png"),
]

FONT = "Noto Sans CJK SC"
SERIF = "Noto Serif CJK SC"
C = {
    "green": RGBColor(0, 176, 80),
    "deep_green": RGBColor(19, 122, 75),
    "pale_green": RGBColor(225, 241, 232),
    "mint": RGBColor(242, 250, 246),
    "black": RGBColor(0, 0, 0),
    "text": RGBColor(28, 33, 38),
    "muted": RGBColor(92, 105, 110),
    "line": RGBColor(203, 219, 211),
    "white": RGBColor(255, 255, 255),
    "orange": RGBColor(180, 80, 20),
    "blue": RGBColor(45, 109, 159),
    "red": RGBColor(178, 63, 54),
    "yellow": RGBColor(247, 183, 49),
}


def ensure_dirs():
    OUT.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)
    RENDERED.mkdir(parents=True, exist_ok=True)


def prepare_background():
    src = SOURCE_SLIDES[0]
    if not src.exists():
        return
    base = Image.open(src).convert("RGB")
    w, h = base.size
    left_src = base.crop((0, 0, int(w * 0.31), h))
    scaled_w = int(left_src.size[0] * 768 / left_src.size[1])
    left = left_src.resize((scaled_w, 768))
    canvas = Image.new("RGB", (1366, 768), "white")
    canvas.paste(left, (0, 0))
    white = Image.new("RGB", canvas.size, "white")
    canvas = Image.blend(canvas, white, 0.38)
    canvas = ImageEnhance.Color(canvas).enhance(0.85)
    fade = Image.new("RGB", canvas.size, "white")
    for x in range(canvas.size[0]):
        alpha = min(1.0, max(0.0, (x - 300) / 220))
        if alpha > 0:
            stripe = Image.blend(canvas.crop((x, 0, x + 1, canvas.size[1])), fade.crop((x, 0, x + 1, canvas.size[1])), alpha)
            canvas.paste(stripe, (x, 0))
    canvas.save(BG_OUT)


def set_run_font(run, size, color="text", bold=False, font=FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = C[color] if isinstance(color, str) else color


def add_text(slide, text, x, y, w, h, size=16, color="text", bold=False, align=None, valign=None, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align if align else PP_ALIGN.LEFT
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = C[color] if isinstance(color, str) else color
    return box


def add_rect(slide, x, y, w, h, fill="white", line="line", radius=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C[fill] if isinstance(fill, str) else fill
    shape.line.color.rgb = C[line] if isinstance(line, str) else line
    shape.line.width = Pt(1)
    return shape


def add_line(slide, x1, y1, x2, y2, color="line", width=1.1, arrow=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = C[color] if isinstance(color, str) else color
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = True
    return line


def add_new_base(prs, title, page, section="方案扩展"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C["white"]
    if BG_OUT.exists():
        slide.shapes.add_picture(str(BG_OUT), 0, 0, width=prs.slide_width, height=prs.slide_height)
    add_rect(slide, 0.5, 0.18, 3.15, 0.42, fill="green", line="green")
    add_text(slide, title, 0.5, 0.25, 3.15, 0.22, size=15, color="black", bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, section, 10.7, 0.32, 1.7, 0.22, size=8.5, color="muted", align=PP_ALIGN.RIGHT)
    add_line(slide, 4.1, 0.62, 12.55, 0.62, color="line", width=0.9)
    add_text(slide, f"{page:02d}", 12.15, 7.03, 0.35, 0.2, size=8.5, color="muted", align=PP_ALIGN.RIGHT)
    return slide


def add_card(slide, x, y, w, h, title, body, accent="green", title_size=13.2, body_size=10.5):
    add_rect(slide, x, y, w, h, fill="white", line="line")
    add_rect(slide, x, y, 0.08, h, fill=accent, line=accent)
    add_text(slide, title, x + 0.2, y + 0.18, w - 0.38, 0.28, size=title_size, color="black", bold=True)
    add_text(slide, body, x + 0.2, y + 0.62, w - 0.38, h - 0.76, size=body_size, color="text")


def add_small_label(slide, x, y, w, text, fill="pale_green"):
    add_rect(slide, x, y, w, 0.36, fill=fill, line="line", radius=True)
    add_text(slide, text, x + 0.08, y + 0.09, w - 0.16, 0.16, size=8.8, color="text", bold=True, align=PP_ALIGN.CENTER)


def add_image(slide, path, x, y, w, h):
    if path.exists():
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y))
        src_ratio = pic.width / pic.height
        dst_ratio = w / h
        if src_ratio > dst_ratio:
            crop = (1 - dst_ratio / src_ratio) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        else:
            crop = (1 - src_ratio / dst_ratio) / 2
            pic.crop_top = crop
            pic.crop_bottom = crop
        pic.left = Inches(x)
        pic.top = Inches(y)
        pic.width = Inches(w)
        pic.height = Inches(h)
    else:
        add_rect(slide, x, y, w, h, fill="mint", line="line")


def add_table_like(slide, x, y, col_widths, row_heights, headers, rows):
    total_w = sum(col_widths)
    total_h = sum(row_heights)
    add_rect(slide, x, y, total_w, total_h, fill="white", line="line")
    cx = x
    for idx, w in enumerate(col_widths):
        add_rect(slide, cx, y, w, row_heights[0], fill="pale_green", line="line")
        add_text(slide, headers[idx], cx + 0.06, y + 0.11, w - 0.12, 0.16, size=9.8, color="black", bold=True, align=PP_ALIGN.CENTER)
        cx += w
    yy = y + row_heights[0]
    for r, row in enumerate(rows):
        cx = x
        for c, cell in enumerate(row):
            add_rect(slide, cx, yy, col_widths[c], row_heights[r + 1], fill="white", line="line")
            align = PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT
            add_text(slide, cell, cx + 0.08, yy + 0.12, col_widths[c] - 0.16, row_heights[r + 1] - 0.18, size=8.8 if c else 9.2, color="text", bold=(c == 0), align=align)
            cx += col_widths[c]
        yy += row_heights[r + 1]


def add_original_slide(prs, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if image_path.exists():
        add_image(slide, image_path, 0, 0, 13.333, 7.5)
    else:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C["white"]
        add_text(slide, f"缺少截图：{image_path}", 0.8, 3.3, 11.5, 0.5, size=20, color="red", align=PP_ALIGN.CENTER)


def build_ppt():
    ensure_dirs()
    prepare_background()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.author = "HanW"
    prs.core_properties.title = "VRX WAM-V 自主航行项目方案扩展版"
    prs.core_properties.subject = "EKF + Dubins + ILOS + State Lattice 方案扩展"

    for slide_img in SOURCE_SLIDES:
        add_original_slide(prs, slide_img)

    slide = add_new_base(prs, "选型结论", 6)
    add_text(slide, "为什么本项目采用 EKF + Dubins + ILOS + State Lattice", 4.25, 0.95, 7.6, 0.45, size=21, color="black", bold=True)
    add_table_like(
        slide,
        4.25,
        1.6,
        [1.45, 2.25, 2.55, 2.35],
        [0.48, 0.92, 0.92, 0.92, 0.92],
        ["模块", "承担问题", "选择原因", "当前边界"],
        [
            ["EKF", "把 GPS/IMU 变成稳定 ENU 状态", "实时、轻量、ROS2 工程成熟", "依赖 GNSS，不解决长期无 GPS"],
            ["Dubins", "生成满足最小转弯半径的几何路径", "解释清晰，可作为基准路径", "只前进，缺少障碍约束"],
            ["ILOS", "沿路径抗横流跟踪", "结构简单，适合欠驱动船", "参数需随速度/扰动调节"],
            ["State Lattice", "障碍和曲率约束下搜索可行路径", "把局部避障转为可验证路径", "复杂场景耗时上升"],
        ],
    )
    add_card(
        slide,
        4.25,
        5.82,
        8.6,
        1.0,
        "一句话结论",
        "EKF 负责“我在哪”，Dubins 负责“粗略怎么转过去”，ILOS 负责“沿着路径怎么走”，State Lattice 负责“有障碍和运动约束时怎么找可执行路径”。",
        "orange",
        body_size=10.0,
    )

    slide = add_new_base(prs, "系统架构", 7)
    add_text(slide, "从传感器到推进器的闭环链路", 4.35, 0.98, 5.2, 0.38, size=20, color="black", bold=True)
    stages = [
        ("GPS / IMU\nLaser / PointCloud", "传感器输入", "green"),
        ("EKF\n状态健康监测", "状态估计", "blue"),
        ("滚动占据栅格\n浮标候选", "局部地图", "deep_green"),
        ("Dubins\nState Lattice A*", "路径规划", "orange"),
        ("ILOS\n曲率限速", "路径跟随", "green"),
        ("PID / 安全层\n左右推进器", "控制执行", "red"),
    ]
    x = 4.2
    for i, (main, label, color) in enumerate(stages):
        add_rect(slide, x, 1.65, 1.12, 1.05, fill="white", line=color)
        add_text(slide, main, x + 0.06, 1.88, 1.0, 0.38, size=8.1, color="text", bold=True, align=PP_ALIGN.CENTER)
        add_small_label(slide, x, 2.86, 1.12, label)
        if i < len(stages) - 1:
            add_line(slide, x + 1.14, 2.16, x + 1.44, 2.16, color="muted", width=1.2, arrow=True)
        x += 1.48
    add_card(slide, 4.25, 4.0, 2.55, 1.45, "观测层", "RViz 显示规划路径、实际轨迹、横向误差、浮标候选、动态船和占据栅格。", "green")
    add_card(slide, 7.0, 4.0, 2.55, 1.45, "监督层", "动态目标不写成静态墙，而是通过 CPA/TCPA 触发 COLREGs 偏置与降速。", "orange")
    add_card(slide, 9.75, 4.0, 2.55, 1.45, "验收层", "自动统计完成率、碰撞、净空、误差、规划耗时和 fallback 次数。", "blue")
    add_text(slide, "架构重点：规划、跟随和安全层分离，任何一层异常都能在状态话题或 RViz 中定位。", 4.25, 6.15, 8.2, 0.36, size=12.5, color="black", bold=True, align=PP_ALIGN.CENTER)

    slide = add_new_base(prs, "状态估计落地", 8)
    add_text(slide, "EKF 不只是公式，而是控制链的统一坐标和速度来源", 4.25, 0.98, 7.7, 0.38, size=19, color="black", bold=True)
    add_card(slide, 4.15, 1.65, 2.55, 2.05, "状态量", "东向 E、北向 N、东速、北速、航向 yaw、艏摇角速度 r，共 6 维状态。", "green")
    add_card(slide, 6.95, 1.65, 2.55, 2.05, "观测融合", "GPS 提供位置和采样速度，IMU 提供航向与角速度，角度差做环绕处理。", "blue")
    add_card(slide, 9.75, 1.65, 2.55, 2.05, "健康判定", "创新门限、协方差、传感器超时共同决定是否使用 robot_localization 或 PlanarEKF 回退。", "orange")
    add_table_like(
        slide,
        4.15,
        4.18,
        [2.0, 2.3, 2.3, 2.15],
        [0.42, 0.68, 0.68, 0.68],
        ["问题", "处理方式", "输出证据", "对控制影响"],
        [
            ["GPS 跳点", "6 sigma 创新门限拒绝", "rejected 数量", "避免目标方向突变"],
            ["IMU 角度跨 ±π", "wrap 角度残差", "yaw 标准差", "避免艏向翻转"],
            ["传感器中断", "超时停车/冻结积分", "status 字段", "先保安全再恢复"],
        ],
    )

    slide = add_new_base(prs, "局部地图", 9)
    add_text(slide, "雷达点云如何变成 State Lattice 可用的障碍约束", 4.25, 0.98, 7.5, 0.38, size=19, color="black", bold=True)
    add_image(slide, ASSETS / "lidar_card.png", 4.15, 1.48, 3.45, 2.48)
    add_card(slide, 7.95, 1.48, 4.05, 0.92, "滚动窗口", "100 m x 100 m，0.5 m 分辨率，随 WAM-V 整格平移。", "green", body_size=9.6)
    add_card(slide, 7.95, 2.62, 4.05, 0.92, "证据更新", "命中提高 log-odds，自由射线降低 log-odds，旧证据随时间衰减。", "blue", body_size=9.6)
    add_card(slide, 7.95, 3.76, 4.05, 0.92, "低浮标确认", "点云滤水、空间聚类、多帧确认后进入浮标候选显示。", "orange", body_size=9.6)
    add_card(slide, 4.15, 5.12, 7.85, 1.02, "为什么不用完整 SLAM 做主线", "VRX Wayfinding 目标由 WGS84 航点发布，开阔水面长期稳定特征少；本项目更需要实时局部障碍记忆，而不是保存全局地图。", "deep_green", body_size=10.2)

    slide = add_new_base(prs, "规划衔接", 10)
    add_text(slide, "Dubins 与 State Lattice 的关系：先快解，再搜索", 4.25, 0.98, 6.8, 0.38, size=19, color="black", bold=True)
    boxes = [
        (4.25, "目标点 + 期望艏向", "官方 Wayfinding 航点或局部 40 m 子目标", "green"),
        (6.35, "解析 Dubins 连接", "无遮挡时快速生成最短曲率受限路径", "blue"),
        (8.45, "碰撞检查失败", "路径穿过占用栅格或膨胀安全圈", "orange"),
        (10.55, "State Lattice A*", "离散位置/航向上搜索可执行路径", "red"),
    ]
    for i, (x, title, body, color) in enumerate(boxes):
        add_card(slide, x, 1.65, 1.75, 2.08, title, body, color, title_size=10.5, body_size=8.4)
        if i < len(boxes) - 1:
            add_line(slide, x + 1.78, 2.68, x + 2.04, 2.68, color="muted", arrow=True)
    add_table_like(
        slide,
        4.25,
        4.35,
        [1.75, 2.35, 2.25, 2.25],
        [0.42, 0.72, 0.72, 0.72],
        ["机制", "触发条件", "输出", "价值"],
        [
            ["子目标", "目标超出滚动地图", "40 m 局部目标", "保持规划在已知区域内"],
            ["重规划", "障碍持续阻断路径", "新 path revision", "绕开新增障碍"],
            ["fallback", "搜索失败或起点异常", "保留旧安全路径", "避免未检查路径覆盖"],
        ],
    )

    slide = add_new_base(prs, "路径跟随", 11)
    add_text(slide, "ILOS 把“路径曲线”变成“期望航向 + 速度约束”", 4.25, 0.98, 7.4, 0.38, size=19, color="black", bold=True)
    add_image(slide, ASSETS / "path_card.png", 4.15, 1.45, 3.55, 2.55)
    add_card(slide, 8.0, 1.45, 4.1, 0.95, "横向误差", "以路径切向为参考，船在路径左侧时 e_y > 0，期望航向向右修正。", "green", body_size=9.4)
    add_card(slide, 8.0, 2.6, 4.1, 0.95, "积分补偿", "b_I 用于抵消稳定横流；避障、制动、低速和误差过大时冻结。", "blue", body_size=9.4)
    add_card(slide, 8.0, 3.75, 4.1, 0.95, "曲率前馈", "路径曲率生成 r_ff，并按横向加速度限制提前降速。", "orange", body_size=9.4)
    add_text(slide, "控制含义：ILOS 不直接躲障，它负责在规划路径已经安全的前提下稳定贴线。", 4.25, 5.55, 7.9, 0.36, size=12.5, color="black", bold=True, align=PP_ALIGN.CENTER)

    slide = add_new_base(prs, "动态避障", 12)
    add_text(slide, "动态船通过预测和规则监督进入控制，不污染静态地图", 4.25, 0.98, 7.8, 0.38, size=19, color="black", bold=True)
    flow = [
        ("目标船检测", "仿真专用 odom", "green"),
        ("速度估计", "alpha-beta 跟踪", "blue"),
        ("风险评估", "CPA / TCPA", "orange"),
        ("会遇分类", "横越 / 对遇 / 追越", "deep_green"),
        ("监督动作", "右转偏置 + 降速", "red"),
    ]
    x = 4.35
    for i, (title, body, color) in enumerate(flow):
        add_card(slide, x, 1.62, 1.42, 1.58, title, body, color, title_size=9.8, body_size=8.2)
        if i < len(flow) - 1:
            add_line(slide, x + 1.44, 2.42, x + 1.72, 2.42, color="muted", arrow=True)
        x += 1.75
    add_table_like(
        slide,
        4.25,
        3.95,
        [1.85, 2.3, 2.3, 2.25],
        [0.42, 0.72, 0.72, 0.72],
        ["场景", "判断依据", "控制策略", "可视化证据"],
        [
            ["横越", "TCPA 在窗口内，DCPA 过小", "向右偏置并降低速度", "风险目标变红"],
            ["对遇", "相对方位接近船艏", "更强右转偏置", "encounter=head_on"],
            ["追越", "后向相对接近", "保守限速/保持距离", "TCPA/DCPA 标签"],
        ],
    )

    slide = add_new_base(prs, "控制输出", 13)
    add_text(slide, "最后一层：把期望速度和艏摇角速度变成左右推进器", 4.25, 0.98, 7.6, 0.38, size=19, color="black", bold=True)
    add_card(slide, 4.2, 1.55, 2.55, 1.55, "速度闭环", "速度 PID 负责巡航、接近目标和弯道限速后的纵向推力。", "green")
    add_card(slide, 6.95, 1.55, 2.55, 1.55, "艏摇闭环", "航向误差 + 曲率前馈生成 r_cmd，并受角速度参考治理器限制。", "blue")
    add_card(slide, 9.7, 1.55, 2.55, 1.55, "推进器分配", "前进推力与差动转向叠加，输出左右推进器指令。", "orange")
    add_card(slide, 4.2, 3.72, 2.55, 1.55, "近障碍安全层", "雷达发现过近障碍时覆盖正常指令，执行制动或安全倒车恢复。", "red")
    add_card(slide, 6.95, 3.72, 2.55, 1.55, "传感器超时", "GPS/IMU/雷达失联时停车，等待有效样本后再恢复动作。", "deep_green")
    add_card(slide, 9.7, 3.72, 2.55, 1.55, "冲突保护", "重复 controller 或重复 Gazebo 会被检查，避免多个发布者抢推进器。", "blue")
    add_text(slide, "控制层目标：宁可保守降速，也不要让规划失败、传感器中断或重复进程直接变成碰撞风险。", 4.25, 6.05, 8.0, 0.36, size=12.2, color="black", bold=True, align=PP_ALIGN.CENTER)

    slide = add_new_base(prs, "实验入口", 14)
    add_text(slide, "从基础闭环到动态避障的推荐演示顺序", 4.25, 0.98, 7.0, 0.38, size=19, color="black", bold=True)
    add_table_like(
        slide,
        4.15,
        1.5,
        [2.2, 3.55, 2.85],
        [0.45, 0.68, 0.68, 0.68, 0.68, 0.68],
        ["场景", "命令", "主要观察点"],
        [
            ["基础 Wayfinding", "simulation.launch.py", "EKF / Dubins / ILOS"],
            ["多航点压力", "multi_waypoint_course.launch.py", "连续重规划和最终艏向"],
            ["固定浮标场", "buoy_course.launch.py", "点云、栅格、浮标候选"],
            ["State Lattice 压力", "lattice_stress.launch.py", "A* 展开、fallback=0"],
            ["COLREGs 学习场", "colregs_learning.launch.py", "动态目标和右转降速"],
        ],
    )
    add_card(
        slide,
        4.15,
        5.55,
        8.6,
        1.12,
        "推荐主入口",
        "ros2 launch han_usv_controller colregs_learning.launch.py：一次启动 Gazebo、RViz、WAM-V、EKF、控制器、12 枚浮标和 2 艘动态目标船。",
        "green",
        body_size=9.0,
    )

    slide = add_new_base(prs, "评价指标", 15)
    add_text(slide, "PPT 最后要用指标收束：不是只会演示，而是可验收", 4.25, 0.98, 7.9, 0.38, size=19, color="black", bold=True)
    metrics = [
        ("100%", "官方任务完成率", "green"),
        ("0", "碰撞次数", "green"),
        ("1.076 m", "官方平均误差", "blue"),
        ("5.952 m", "最小净空", "orange"),
        ("0", "State Lattice fallback", "green"),
        ("189", "COLREGs 激活样本", "red"),
    ]
    x_positions = [4.25, 6.85, 9.45, 4.25, 6.85, 9.45]
    y_positions = [1.62, 1.62, 1.62, 3.15, 3.15, 3.15]
    for (value, label, color), x, y in zip(metrics, x_positions, y_positions):
        add_rect(slide, x, y, 2.25, 1.08, fill="white", line="line")
        add_text(slide, value, x + 0.1, y + 0.2, 2.05, 0.3, size=19, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.1, y + 0.66, 2.05, 0.2, size=9.4, color="muted", align=PP_ALIGN.CENTER)
    add_card(slide, 4.25, 5.0, 8.05, 0.9, "自动验收命令", "bash han_usv_controller/scripts/run_colregs_learning.sh", "green", body_size=12.0)
    add_text(slide, "展示建议：先给 RViz 画面，再给 aggregate.json 指标，这样“能看见”和“能量化”都成立。", 4.25, 6.22, 8.05, 0.3, size=11.8, color="black", bold=True, align=PP_ALIGN.CENTER)

    slide = add_new_base(prs, "边界与后续", 16)
    add_text(slide, "把边界说清楚，反而能让方案显得更工程化", 4.25, 0.98, 7.0, 0.38, size=19, color="black", bold=True)
    add_card(slide, 4.15, 1.55, 2.55, 2.0, "当前边界", "动态船检测来自仿真专用 Topic；激光浮标候选不能识别红/绿语义；COLREGs 是学习监督层。", "red")
    add_card(slide, 6.95, 1.55, 2.55, 2.0, "近期扩展", "接入视觉颜色识别；把动态目标检测从理想 odom 换成雷达/视觉融合；补更多会遇场景。", "orange")
    add_card(slide, 9.75, 1.55, 2.55, 2.0, "研究扩展", "Factor Graph/iSAM2 用于离线轨迹优化；MPC 用于统一处理动力学、速度和避障约束。", "blue")
    add_table_like(
        slide,
        4.15,
        4.2,
        [2.2, 3.2, 3.0],
        [0.42, 0.72, 0.72],
        ["方向", "进入条件", "预期收益"],
        [
            ["SLAM/图优化", "港池、岸线特征稳定、GNSS 不可靠", "提高长期定位一致性"],
            ["MPC", "需要同时优化轨迹、速度、避碰和能耗", "提升复杂约束下的平滑性"],
        ],
    )

    slide = add_new_base(prs, "总结", 17)
    add_text(slide, "本项目的汇报主线可以收束为三句话", 4.25, 0.98, 6.0, 0.38, size=20, color="black", bold=True)
    add_card(slide, 4.35, 1.78, 7.65, 0.95, "1. 方案不是堆算法", "EKF、Dubins、ILOS、State Lattice 分别对应定位、粗规划、路径跟随和约束规划，职责清晰。", "green", body_size=10.5)
    add_card(slide, 4.35, 3.08, 7.65, 0.95, "2. 系统不是只会跑", "RViz、status 和自动评测把每次决策、每次风险、每次失败都变成可检查证据。", "blue", body_size=10.5)
    add_card(slide, 4.35, 4.38, 7.65, 0.95, "3. 后续不是空泛扩展", "视觉语义、真实目标检测、图优化和 MPC 都能沿着现有接口逐步接入。", "orange", body_size=10.5)
    add_text(slide, "最终表达：这是一套面向 VRX 学习、算法对比和后续研究的可复现自主航行实验平台。", 4.35, 6.05, 7.65, 0.38, size=13.8, color="black", bold=True, align=PP_ALIGN.CENTER)

    prs.save(PPTX_OUT)


def write_notes():
    NOTES_OUT.write_text(
        dedent(
            """\
            # VRX WAM-V 自主航行项目方案扩展版讲稿

            ## 1-5. 原有方案对比页
            这五页先完成方案横向比较：状态估计、曲线生成、路径跟随和运动规划分别有哪些常见方法。重点不要逐字念表格，而是强调我们最终选择的是一条轻量、可解释、可实时运行的工程链。

            ## 6. 选型结论
            EKF 解决“我在哪”，Dubins 解决“粗略怎么转过去”，ILOS 解决“沿路径怎么走”，State Lattice 解决“有障碍和运动约束时怎么规划可执行路径”。这页要把前面四张对比表收束成项目选型。

            ## 7. 系统架构
            从传感器到推进器的主链路是 GPS/IMU/雷达输入，经过 EKF、滚动地图、Dubins/State Lattice、ILOS、PID 和安全层，最后输出左右推进器。旁路还有 RViz 可视化、COLREGs 监督和自动验收。

            ## 8. 状态估计落地
            这里讲 EKF 的工程作用：统一 ENU 坐标、提供速度和航向、做健康检测。它不是为了展示滤波公式，而是为了让规划和控制有稳定、连续、可检查的状态输入。

            ## 9. 局部地图
            雷达和点云进入 100 m 滚动占据栅格，旧证据会衰减，低浮标通过滤水和聚类变成黄色候选。强调本项目不把 SLAM 当主线，是因为 Wayfinding 目标来自 GNSS 航点，当前更需要实时局部障碍记忆。

            ## 10. 规划衔接
            无障碍时优先使用解析 Dubins，速度快且解释清楚；碰撞检查失败时进入 State Lattice A*，在离散位置和航向上搜索可执行路径。这样能兼顾计算效率和障碍约束。

            ## 11. 路径跟随
            ILOS 把路径变成期望航向，并通过积分项补偿稳定横流。避障、制动、低速和大误差时积分冻结，避免把临时动作学成长期偏置。曲率前馈和限速让船不会高速硬闯小半径圆弧。

            ## 12. 动态避障
            动态船先估计速度，再计算 CPA/TCPA 和会遇类型，最后给 ILOS/PID 施加右转偏置和降速。动态目标不直接写入静态地图，避免把移动船当成固定墙。

            ## 13. 控制输出
            控制层负责把期望速度和艏摇角速度分配到左右推进器。安全层会覆盖正常控制，处理近障碍、传感器超时和重复进程等风险。

            ## 14. 实验入口
            推荐按基础 Wayfinding、多航点、固定浮标、State Lattice 压力、COLREGs 学习场逐步演示。主入口是 `ros2 launch han_usv_controller colregs_learning.launch.py`。

            ## 15. 评价指标
            汇报不要只停留在截图，最后用自动评测指标证明：完成率 100%、碰撞 0、平均误差 1.076 m、最小净空 5.952 m、fallback 0、COLREGs 激活 189 个样本。

            ## 16. 边界与后续
            当前动态船检测来自仿真专用 Topic，激光不能识别红绿语义，COLREGs 只是学习监督层。后续可以接入视觉语义、真实目标检测、图优化和 MPC。

            ## 17. 总结
            收束为三句话：方案不是堆算法，系统不是只会跑，后续不是空泛扩展。最终定位是一套面向 VRX 学习、算法对比和后续研究的可复现自主航行实验平台。
            """
        ),
        encoding="utf-8",
    )


if __name__ == "__main__":
    build_ppt()
    write_notes()
    print(PPTX_OUT)
    print(NOTES_OUT)
